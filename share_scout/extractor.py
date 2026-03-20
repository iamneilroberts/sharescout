"""Extract text content from files and compute content hashes."""

import hashlib
import logging
from dataclasses import dataclass, field
from pathlib import Path

logger = logging.getLogger(__name__)


# --- Structured extraction dataclasses ---

@dataclass
class TextSection:
    text: str
    position: str       # "beginning", "middle", "end", "heading:<title>", "sheet:<name>", "slide:<n>"
    char_offset: int


@dataclass
class ExtractedImage:
    data: bytes
    format: str         # "png", "jpeg"
    width: int
    height: int
    content_hash: str   # SHA-256 for dedup
    char_offset: int    # position in document
    source: str         # "pdf_page_3", "docx_para_12", "pptx_slide_5"


@dataclass
class DocumentContent:
    sections: list[TextSection] = field(default_factory=list)
    images: list[ExtractedImage] = field(default_factory=list)
    total_chars: int = 0
    page_count: int | None = None
    sheet_count: int | None = None
    slide_count: int | None = None
    heading_count: int = 0
    structure_type: str = "flat"  # "flat", "headed", "paged", "tabular", "slides"


# --- Original text extraction (backward compat) ---

def extract_text(filepath: str, max_chars: int = 4000) -> str | None:
    """Extract text from a file based on its extension.

    Returns up to max_chars of text, or None if extraction fails.
    """
    # Skip Office temp/lock files (start with ~$)
    filename = Path(filepath).name
    if filename.startswith("~$"):
        return None

    ext = Path(filepath).suffix.lower()

    try:
        if ext in (".txt", ".md", ".mdx", ".csv", ".html", ".rtf", ".log",
                   ".rst", ".adoc", ".yaml", ".yml", ".toml", ".json",
                   ".py", ".ts", ".tsx", ".js", ".jsx", ".sql", ".css", ".sh",
                   ".ps1", ".psm1", ".psd1", ".bat", ".cmd", ".vbs", ".rb",
                   ".go", ".rs", ".c", ".cpp", ".h", ".cs", ".java", ".xml"):
            return _extract_plain_text(filepath, max_chars)
        elif ext == ".pdf":
            return _extract_pdf(filepath, max_chars)
        elif ext in (".docx", ".doc"):
            return _extract_docx(filepath, max_chars)
        elif ext == ".xlsx":
            return _extract_xlsx(filepath, max_chars)
        elif ext == ".xls":
            return _extract_xls(filepath, max_chars)
        elif ext == ".pptx":
            return _extract_pptx(filepath, max_chars)
        else:
            logger.debug("No extractor for extension: %s", ext)
            return None
    except Exception as e:
        logger.warning("Extraction failed for %s: %s", filepath, e)
        return None


def compute_partial_hash(filepath: str, hash_bytes: int = 65536) -> str | None:
    """Compute SHA-256 hash of the first hash_bytes of a file."""
    try:
        with open(filepath, "rb") as f:
            data = f.read(hash_bytes)
        return hashlib.sha256(data).hexdigest()
    except (OSError, PermissionError) as e:
        logger.warning("Cannot hash %s: %s", filepath, e)
        return None


# --- Structured extraction ---

def extract_structured(filepath: str, preset_config: dict = None) -> DocumentContent | None:
    """Extract structured content from a file, returning sections and images.

    Returns DocumentContent with text sections and extracted images,
    or None if extraction fails entirely.
    """
    filename = Path(filepath).name
    if filename.startswith("~$"):
        return None

    ext = Path(filepath).suffix.lower()
    preset_config = preset_config or {}
    extraction_cfg = preset_config.get("extraction", {})

    try:
        if ext in (".txt", ".md", ".mdx", ".csv", ".html", ".rtf", ".log",
                   ".rst", ".adoc", ".yaml", ".yml", ".toml", ".json",
                   ".py", ".ts", ".tsx", ".js", ".jsx", ".sql", ".css", ".sh",
                   ".ps1", ".psm1", ".psd1", ".bat", ".cmd", ".vbs", ".rb",
                   ".go", ".rs", ".c", ".cpp", ".h", ".cs", ".java", ".xml"):
            return _extract_structured_plain_text(filepath)
        elif ext == ".pdf":
            return _extract_structured_pdf(filepath)
        elif ext in (".docx", ".doc"):
            return _extract_structured_docx(filepath)
        elif ext == ".xlsx":
            max_rows = extraction_cfg.get("xlsx_max_rows_per_sheet", 50)
            return _extract_structured_xlsx(filepath, max_rows)
        elif ext == ".xls":
            max_rows = extraction_cfg.get("xlsx_max_rows_per_sheet", 50)
            return _extract_structured_xls(filepath, max_rows)
        elif ext == ".pptx":
            max_slides = extraction_cfg.get("pptx_max_slides", 20)
            return _extract_structured_pptx(filepath, max_slides)
        else:
            logger.debug("No structured extractor for extension: %s", ext)
            return None
    except Exception as e:
        logger.warning("Structured extraction failed for %s: %s", filepath, e)
        return None


# --- Plain text extractors ---

def _extract_plain_text(filepath: str, max_chars: int) -> str | None:
    """Read plain text with encoding detection."""
    try:
        with open(filepath, "r", encoding="utf-8") as f:
            return f.read(max_chars)
    except UnicodeDecodeError:
        pass

    try:
        import chardet
        with open(filepath, "rb") as f:
            raw = f.read(max_chars * 2)
        detected = chardet.detect(raw)
        encoding = detected.get("encoding", "latin-1") or "latin-1"
        return raw.decode(encoding, errors="replace")[:max_chars]
    except Exception as e:
        logger.warning("Plain text extraction failed for %s: %s", filepath, e)
        return None


def _read_full_text(filepath: str) -> str | None:
    """Read the full file text with encoding detection."""
    try:
        with open(filepath, "r", encoding="utf-8") as f:
            return f.read()
    except UnicodeDecodeError:
        pass
    try:
        import chardet
        with open(filepath, "rb") as f:
            raw = f.read()
        detected = chardet.detect(raw[:65536])
        encoding = detected.get("encoding", "latin-1") or "latin-1"
        return raw.decode(encoding, errors="replace")
    except Exception as e:
        logger.warning("Full text read failed for %s: %s", filepath, e)
        return None


def _extract_structured_plain_text(filepath: str) -> DocumentContent | None:
    """Extract structured sections from plain text / markdown."""
    text = _read_full_text(filepath)
    if not text or not text.strip():
        return None

    ext = Path(filepath).suffix.lower()
    sections = []
    offset = 0

    # Split at markdown headings if present
    if ext in (".md", ".mdx", ".rst", ".adoc") or text.lstrip().startswith("#"):
        import re
        heading_pattern = re.compile(r'^(#{1,6})\s+(.+)$', re.MULTILINE)
        matches = list(heading_pattern.finditer(text))

        if matches:
            # Text before first heading
            if matches[0].start() > 0:
                pre_text = text[:matches[0].start()].strip()
                if pre_text:
                    sections.append(TextSection(text=pre_text, position="beginning", char_offset=0))

            for i, match in enumerate(matches):
                heading_title = match.group(2).strip()
                start = match.start()
                end = matches[i + 1].start() if i + 1 < len(matches) else len(text)
                section_text = text[start:end].strip()
                if section_text:
                    sections.append(TextSection(
                        text=section_text,
                        position=f"heading:{heading_title}",
                        char_offset=start,
                    ))

            return DocumentContent(
                sections=sections,
                total_chars=len(text),
                heading_count=len(matches),
                structure_type="headed",
            )

    # No headings — split into beginning/middle/end for large files
    if len(text) > 8000:
        third = len(text) // 3
        sections = [
            TextSection(text=text[:third], position="beginning", char_offset=0),
            TextSection(text=text[third:2*third], position="middle", char_offset=third),
            TextSection(text=text[2*third:], position="end", char_offset=2*third),
        ]
        return DocumentContent(sections=sections, total_chars=len(text), structure_type="flat")

    # Small file — single section
    sections = [TextSection(text=text, position="beginning", char_offset=0)]
    return DocumentContent(sections=sections, total_chars=len(text), structure_type="flat")


# --- PDF extractors ---

def _extract_pdf(filepath: str, max_chars: int) -> str | None:
    """Extract text from PDF using pymupdf."""
    import pymupdf

    text_parts = []
    total = 0
    with pymupdf.open(filepath) as doc:
        for page in doc:
            page_text = page.get_text()
            text_parts.append(page_text)
            total += len(page_text)
            if total >= max_chars:
                break

    text = "\n".join(text_parts)[:max_chars]
    return text if text.strip() else None


def _extract_structured_pdf(filepath: str) -> DocumentContent | None:
    """Extract structured sections and images from PDF."""
    import pymupdf

    sections = []
    images = []
    total_chars = 0

    with pymupdf.open(filepath) as doc:
        page_count = len(doc)

        for page_idx, page in enumerate(doc):
            page_text = page.get_text()
            if page_text.strip():
                sections.append(TextSection(
                    text=page_text,
                    position=f"page:{page_idx + 1}",
                    char_offset=total_chars,
                ))
                total_chars += len(page_text)

            # Extract images
            try:
                for img_info in page.get_images(full=True):
                    xref = img_info[0]
                    try:
                        img_data = doc.extract_image(xref)
                        if not img_data:
                            continue
                        width = img_data.get("width", 0)
                        height = img_data.get("height", 0)
                        # Skip tiny images (logos, icons, spacers)
                        if width < 100 or height < 100:
                            continue
                        raw = img_data["image"]
                        content_hash = hashlib.sha256(raw).hexdigest()
                        fmt = img_data.get("ext", "png")
                        images.append(ExtractedImage(
                            data=raw,
                            format=fmt,
                            width=width,
                            height=height,
                            content_hash=content_hash,
                            char_offset=total_chars,
                            source=f"pdf_page_{page_idx + 1}",
                        ))
                    except Exception:
                        continue
            except Exception:
                pass

    if not sections:
        return None

    return DocumentContent(
        sections=sections,
        images=images,
        total_chars=total_chars,
        page_count=page_count,
        structure_type="paged",
    )


# --- DOCX extractors ---

def _extract_docx(filepath: str, max_chars: int) -> str | None:
    """Extract text from .docx using python-docx."""
    from docx import Document

    doc = Document(filepath)
    text_parts = []
    total = 0
    for para in doc.paragraphs:
        text_parts.append(para.text)
        total += len(para.text)
        if total >= max_chars:
            break

    text = "\n".join(text_parts)[:max_chars]
    return text if text.strip() else None


def _extract_structured_docx(filepath: str) -> DocumentContent | None:
    """Extract structured sections and images from DOCX."""
    from docx import Document

    doc = Document(filepath)
    sections = []
    images = []
    total_chars = 0
    heading_count = 0

    # Group paragraphs by headings
    current_section_parts = []
    current_heading = None
    section_start_offset = 0

    for para_idx, para in enumerate(doc.paragraphs):
        is_heading = para.style and para.style.name and para.style.name.startswith("Heading")

        if is_heading and current_section_parts:
            # Save previous section
            section_text = "\n".join(current_section_parts)
            position = f"heading:{current_heading}" if current_heading else "beginning"
            sections.append(TextSection(
                text=section_text,
                position=position,
                char_offset=section_start_offset,
            ))
            current_section_parts = []
            section_start_offset = total_chars

        if is_heading:
            heading_count += 1
            current_heading = para.text.strip()

        current_section_parts.append(para.text)
        total_chars += len(para.text) + 1  # +1 for newline

    # Save last section
    if current_section_parts:
        section_text = "\n".join(current_section_parts)
        position = f"heading:{current_heading}" if current_heading else "beginning"
        sections.append(TextSection(
            text=section_text,
            position=position,
            char_offset=section_start_offset,
        ))

    # Extract images from relationships
    try:
        from docx.opc.constants import RELATIONSHIP_TYPE as RT
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                try:
                    image_part = rel.target_part
                    raw = image_part.blob
                    content_hash = hashlib.sha256(raw).hexdigest()
                    # Determine format from content type
                    ct = image_part.content_type or ""
                    fmt = "png"
                    if "jpeg" in ct or "jpg" in ct:
                        fmt = "jpeg"
                    elif "png" in ct:
                        fmt = "png"
                    # We don't have easy access to dimensions without PIL,
                    # so we'll set them to 0 and let the analyzer handle it
                    images.append(ExtractedImage(
                        data=raw,
                        format=fmt,
                        width=0,
                        height=0,
                        content_hash=content_hash,
                        char_offset=0,
                        source=f"docx_image",
                    ))
                except Exception:
                    continue
    except Exception:
        pass

    if not sections:
        return None

    structure_type = "headed" if heading_count > 0 else "flat"

    return DocumentContent(
        sections=sections,
        images=images,
        total_chars=total_chars,
        heading_count=heading_count,
        structure_type=structure_type,
    )


# --- XLSX extractors ---

def _extract_xlsx(filepath: str, max_chars: int) -> str | None:
    """Extract sheet names and header rows from .xlsx using openpyxl."""
    from openpyxl import load_workbook

    wb = load_workbook(filepath, read_only=True, data_only=True)
    text_parts = []
    total = 0

    for sheet_name in wb.sheetnames:
        text_parts.append(f"Sheet: {sheet_name}")
        total += len(sheet_name) + 8
        ws = wb[sheet_name]
        for row in ws.iter_rows(max_row=20, values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            line = " | ".join(cells)
            text_parts.append(line)
            total += len(line)
            if total >= max_chars:
                break
        if total >= max_chars:
            break

    wb.close()
    text = "\n".join(text_parts)[:max_chars]
    return text if text.strip() else None


def _extract_structured_xlsx(filepath: str, max_rows: int = 50) -> DocumentContent | None:
    """Extract structured sections from XLSX — one section per sheet."""
    from openpyxl import load_workbook

    wb = load_workbook(filepath, read_only=True, data_only=True)
    sections = []
    total_chars = 0

    for sheet_name in wb.sheetnames:
        text_parts = [f"Sheet: {sheet_name}"]
        ws = wb[sheet_name]
        row_count = 0
        for row in ws.iter_rows(max_row=max_rows, values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            line = " | ".join(cells)
            text_parts.append(line)
            row_count += 1

        section_text = "\n".join(text_parts)
        sections.append(TextSection(
            text=section_text,
            position=f"sheet:{sheet_name}",
            char_offset=total_chars,
        ))
        total_chars += len(section_text)

    wb.close()

    if not sections:
        return None

    return DocumentContent(
        sections=sections,
        total_chars=total_chars,
        sheet_count=len(sections),
        structure_type="tabular",
    )


# --- XLS extractor ---

def _extract_xls(filepath: str, max_chars: int) -> str | None:
    """Extract text from legacy .xls using xlrd."""
    try:
        import xlrd
    except ImportError:
        logger.debug("xlrd not installed, cannot read .xls files. Run: pip install xlrd")
        return None

    try:
        wb = xlrd.open_workbook(filepath)
        text_parts = []
        total = 0
        for sheet in wb.sheets():
            text_parts.append(f"Sheet: {sheet.name}")
            total += len(sheet.name) + 8
            for row_idx in range(min(sheet.nrows, 20)):
                cells = [str(sheet.cell_value(row_idx, col)) for col in range(sheet.ncols)]
                line = " | ".join(c for c in cells if c.strip())
                if line:
                    text_parts.append(line)
                    total += len(line)
                if total >= max_chars:
                    break
            if total >= max_chars:
                break
        text = "\n".join(text_parts)[:max_chars]
        return text if text.strip() else None
    except Exception as e:
        logger.warning("XLS extraction failed for %s: %s", filepath, e)
        return None


def _extract_structured_xls(filepath: str, max_rows: int = 50) -> DocumentContent | None:
    """Extract structured sections from legacy .xls — one section per sheet."""
    try:
        import xlrd
    except ImportError:
        logger.debug("xlrd not installed, cannot read .xls files")
        return None

    try:
        wb = xlrd.open_workbook(filepath)
        sections = []
        total_chars = 0

        for sheet in wb.sheets():
            text_parts = [f"Sheet: {sheet.name}"]
            for row_idx in range(min(sheet.nrows, max_rows)):
                cells = [str(sheet.cell_value(row_idx, col)) for col in range(sheet.ncols)]
                line = " | ".join(c for c in cells if c.strip())
                if line:
                    text_parts.append(line)

            section_text = "\n".join(text_parts)
            sections.append(TextSection(
                text=section_text,
                position=f"sheet:{sheet.name}",
                char_offset=total_chars,
            ))
            total_chars += len(section_text)

        if not sections:
            return None

        return DocumentContent(
            sections=sections,
            total_chars=total_chars,
            sheet_count=len(sections),
            structure_type="tabular",
        )
    except Exception as e:
        logger.warning("Structured XLS extraction failed for %s: %s", filepath, e)
        return None


# --- PPTX extractors ---

def _extract_pptx(filepath: str, max_chars: int) -> str | None:
    """Extract slide text from .pptx using python-pptx."""
    try:
        from pptx import Presentation
    except ImportError:
        logger.debug("python-pptx not installed, skipping PPTX extraction")
        return None

    prs = Presentation(filepath)
    text_parts = []
    total = 0

    for i, slide in enumerate(prs.slides):
        if i >= 10:
            break
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text_parts.append(shape.text)
                total += len(shape.text)
                if total >= max_chars:
                    break
        if total >= max_chars:
            break

    text = "\n".join(text_parts)[:max_chars]
    return text if text.strip() else None


def _extract_structured_pptx(filepath: str, max_slides: int = 20) -> DocumentContent | None:
    """Extract structured sections and images from PPTX — one section per slide."""
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        logger.debug("python-pptx not installed, skipping PPTX extraction")
        return None

    prs = Presentation(filepath)
    sections = []
    images = []
    total_chars = 0
    slide_count = len(prs.slides)

    for i, slide in enumerate(prs.slides):
        if i >= max_slides:
            break

        text_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text_parts.append(shape.text)

            # Extract images
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image = shape.image
                    raw = image.blob
                    content_hash = hashlib.sha256(raw).hexdigest()
                    fmt = image.content_type.split("/")[-1] if image.content_type else "png"
                    if fmt == "jpeg":
                        fmt = "jpeg"
                    elif fmt not in ("png", "jpeg"):
                        fmt = "png"
                    images.append(ExtractedImage(
                        data=raw,
                        format=fmt,
                        width=shape.width or 0,
                        height=shape.height or 0,
                        content_hash=content_hash,
                        char_offset=total_chars,
                        source=f"pptx_slide_{i + 1}",
                    ))
            except Exception:
                continue

        section_text = "\n".join(text_parts)
        if section_text.strip():
            sections.append(TextSection(
                text=section_text,
                position=f"slide:{i + 1}",
                char_offset=total_chars,
            ))
            total_chars += len(section_text)

    if not sections:
        return None

    return DocumentContent(
        sections=sections,
        images=images,
        total_chars=total_chars,
        slide_count=slide_count,
        structure_type="slides",
    )
